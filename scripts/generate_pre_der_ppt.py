"""
Generate per-opportunity PowerPoint presentation from Pre-DER Agent results.
Output: output/pre_der_agent/presentation_pre_der_exp1.pptx
"""
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import lxml.etree as etree

ROOT = Path(__file__).parent.parent
JSON_PATH = ROOT / "output/pre_der_agent/results_pre_der_exp1_extracted.json"
OUT_PATH = ROOT / "output/pre_der_agent/presentation_pre_der_exp1.pptx"

# ── Colors ────────────────────────────────────────────────────────────────────
C_HEADER_BG   = RGBColor(0x1F, 0x38, 0x64)  # dark navy
C_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_GRAY  = RGBColor(0xF2, 0xF2, 0xF2)
C_MID_GRAY    = RGBColor(0xBF, 0xBF, 0xBF)
C_DARK_GRAY   = RGBColor(0x40, 0x40, 0x40)
C_SECTION     = RGBColor(0x26, 0x26, 0x26)
C_LEFT_BG     = RGBColor(0xFA, 0xFA, 0xFA)
C_INPUT_BG    = RGBColor(0xED, 0xED, 0xED)
C_DIVIDER     = RGBColor(0xD9, 0xD9, 0xD9)

BG_COLORS = {
    "IDG": RGBColor(0x44, 0x72, 0xC4),
    "DCG": RGBColor(0x70, 0x30, 0xA0),
    "SSG": RGBColor(0x00, 0xB0, 0xF0),
}
RANK_COLORS = [
    RGBColor(0xFF, 0xC0, 0x00),  # #1 gold
    RGBColor(0x80, 0x80, 0x80),  # #2 silver
    RGBColor(0xC5, 0x5A, 0x11),  # #3 bronze
]
CONF_COLORS = {
    "High":   RGBColor(0x70, 0xAD, 0x47),
    "Medium": RGBColor(0xFF, 0xC0, 0x00),
    "Low":    RGBColor(0xFF, 0x00, 0x00),
}

# ── Geometry (inches) ─────────────────────────────────────────────────────────
W_SLIDE, H_SLIDE = 13.33, 7.5
H_HEADER = 1.0
H_FOOTER = 0.38
PAD = 0.18
W_LEFT = 4.9
X_RIGHT = W_LEFT + PAD * 2
W_RIGHT = W_SLIDE - X_RIGHT - PAD
Y_BODY = H_HEADER + 0.12
H_BODY = H_SLIDE - H_HEADER - H_FOOTER - 0.12


def rgb(r, g, b):
    return RGBColor(r, g, b)


def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width_pt=0.75):
    from pptx.util import Pt as _Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.line.color.rgb = line_color if line_color else (fill_color or C_WHITE)
    if line_color is None:
        shape.line.fill.background()
        shape.line.width = 0
    else:
        shape.line.color.rgb = line_color
        from pptx.util import Pt as _Pt
        shape.line.width = _Pt(line_width_pt)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape


def add_textbox(slide, x, y, w, h, text, font_size=10, bold=False, italic=False,
                color=C_DARK_GRAY, align=PP_ALIGN.LEFT, word_wrap=True,
                fill_color=None, line_color=None):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = word_wrap
    if fill_color:
        txb.fill.solid()
        txb.fill.fore_color.rgb = fill_color
    else:
        txb.fill.background()
    if line_color:
        txb.line.color.rgb = line_color
        txb.line.width = Pt(0.5)
    else:
        txb.line.fill.background()

    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def set_tf_margin(txb, left=0.05, top=0.03, right=0.05, bottom=0.03):
    """Set text frame internal margins in inches."""
    tf = txb.text_frame
    tf.margin_left  = Inches(left)
    tf.margin_right = Inches(right)
    tf.margin_top   = Inches(top)
    tf.margin_bottom = Inches(bottom)


def add_label(slide, x, y, w, text):
    """Section label — small caps style, dark, bold."""
    txb = add_textbox(slide, x, y, w, 0.22, text,
                      font_size=7.5, bold=True, color=C_SECTION)
    set_tf_margin(txb, 0, 0, 0, 0)
    return txb


def add_multiline_textbox(slide, x, y, w, h, text, font_size=8.5,
                          color=C_DARK_GRAY, italic=False, fill_color=None, line_color=None):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = True
    set_tf_margin(txb, 0.1, 0.08, 0.1, 0.05)
    if fill_color:
        txb.fill.solid()
        txb.fill.fore_color.rgb = fill_color
    else:
        txb.fill.background()
    if line_color:
        txb.line.color.rgb = line_color
        txb.line.width = Pt(0.5)
    else:
        txb.line.fill.background()
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_slide(prs, record):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    opp_id    = record.get("id", "?")
    title     = record.get("title") or "Untitled"
    bg        = record.get("bg", "").upper()
    raw_text  = record.get("raw_text") or ""
    desc      = record.get("description") or "N/A"
    topk      = record.get("topk") or []
    status    = record.get("status", "unknown")
    provider  = record.get("provider_used", "unknown")

    # Truncate long strings
    if len(title) > 85:
        title = title[:82] + "…"
    if len(raw_text) > 480:
        raw_text = raw_text[:477] + "…"

    # ── Header background ─────────────────────────────────────────────────────
    add_rect(slide, 0, 0, W_SLIDE, H_HEADER, fill_color=C_HEADER_BG)

    # Header: ID + Title
    header_text = f"#{opp_id}  {title}"
    txb = add_textbox(slide, PAD, 0.1, W_SLIDE - 2.0, 0.8, header_text,
                      font_size=18, bold=True, color=C_WHITE)
    set_tf_margin(txb, 0.05, 0.05)

    # BG badge
    bg_color = BG_COLORS.get(bg, RGBColor(0x70, 0x70, 0x70))
    badge = add_rect(slide, W_SLIDE - 1.55, 0.22, 1.3, 0.52, fill_color=bg_color)
    tf = badge.text_frame
    tf.word_wrap = False
    set_tf_margin(badge, 0.05, 0.08)
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = bg if bg else "N/A"
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = C_WHITE

    # ── Left panel background ─────────────────────────────────────────────────
    add_rect(slide, 0, H_HEADER, W_LEFT + PAD, H_SLIDE - H_HEADER, fill_color=C_LEFT_BG)

    # Section: SALES INPUT
    y_input = Y_BODY
    add_label(slide, PAD, y_input, W_LEFT, "SALES INPUT")
    y_input += 0.22

    h_input = 2.35
    add_multiline_textbox(slide, PAD, y_input, W_LEFT, h_input,
                          f'"{raw_text}"',
                          font_size=8, italic=True, color=RGBColor(0x40,0x40,0x40),
                          fill_color=C_INPUT_BG, line_color=C_MID_GRAY)

    # Section: AI EXTRACTED DESCRIPTION
    y_desc = y_input + h_input + 0.18
    add_label(slide, PAD, y_desc, W_LEFT, "AI EXTRACTED DESCRIPTION")
    y_desc += 0.22

    h_desc = H_SLIDE - H_FOOTER - y_desc - 0.1
    add_multiline_textbox(slide, PAD, y_desc, W_LEFT, h_desc,
                          desc, font_size=9, color=C_DARK_GRAY)

    # ── Vertical divider ──────────────────────────────────────────────────────
    div_x = W_LEFT + PAD
    add_rect(slide, div_x, H_HEADER + 0.05,
             0.02, H_SLIDE - H_HEADER - H_FOOTER - 0.1,
             fill_color=C_DIVIDER)

    # ── Right panel: Recommendations ──────────────────────────────────────────
    x_r = div_x + 0.1
    y_r = Y_BODY
    add_label(slide, x_r, y_r, W_RIGHT, "TOP RECOMMENDED PRODUCTS")
    y_r += 0.25

    if not topk:
        add_textbox(slide, x_r, y_r, W_RIGHT, 0.5,
                    "No recommendations available (extraction failed or no match).",
                    font_size=9, color=RGBColor(0xFF,0x00,0x00))
    else:
        # Available height for cards
        avail_h = H_SLIDE - H_FOOTER - y_r - 0.1
        n = min(len(topk), 3)
        card_gap = 0.12
        card_h = (avail_h - card_gap * (n - 1)) / n

        for i, item in enumerate(topk[:3]):
            rank_color = RANK_COLORS[i]
            name       = item.get("name", "Unknown")
            path_str   = item.get("path_str", "")
            score      = item.get("score", 0.0)
            pn_count   = item.get("pn_count", 0)
            conf       = item.get("level_label", "")
            conf_color = CONF_COLORS.get(conf, C_DARK_GRAY)

            y_card = y_r + i * (card_h + card_gap)

            # Card background
            add_rect(slide, x_r, y_card, W_RIGHT, card_h,
                     fill_color=C_WHITE, line_color=C_MID_GRAY, line_width_pt=0.5)

            # Left accent bar
            add_rect(slide, x_r, y_card, 0.08, card_h, fill_color=rank_color)

            inner_x = x_r + 0.14
            inner_w = W_RIGHT - 0.22

            # Rank label
            rank_txb = add_textbox(slide, inner_x, y_card + 0.04, 0.6, 0.28,
                                   f"#{i+1}", font_size=13, bold=True,
                                   color=rank_color)
            set_tf_margin(rank_txb, 0, 0)

            # Product name
            name_txb = add_textbox(slide, inner_x + 0.55, y_card + 0.05,
                                   inner_w - 0.55, 0.32,
                                   name, font_size=11.5, bold=True, color=C_DARK_GRAY)
            set_tf_margin(name_txb, 0, 0)

            # Path
            path_txb = add_textbox(slide, inner_x, y_card + 0.35,
                                   inner_w, 0.28,
                                   path_str, font_size=7.5, italic=True,
                                   color=RGBColor(0x80,0x80,0x80))
            set_tf_margin(path_txb, 0, 0)

            # Score line
            score_y = y_card + card_h - 0.30
            meta = f"Score: {score:.2f}   |   {pn_count} PNs"
            meta_txb = add_textbox(slide, inner_x, score_y, inner_w - 1.2, 0.25,
                                   meta, font_size=8, color=C_DARK_GRAY)
            set_tf_margin(meta_txb, 0, 0)

            # Confidence badge
            conf_x = x_r + W_RIGHT - 1.25
            conf_badge = add_rect(slide, conf_x, score_y - 0.02, 1.1, 0.26,
                                  fill_color=conf_color)
            tf = conf_badge.text_frame
            tf.word_wrap = False
            set_tf_margin(conf_badge, 0.04, 0.02)
            para = tf.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            run = para.add_run()
            run.text = conf
            run.font.size = Pt(8)
            run.font.bold = True
            run.font.color.rgb = C_WHITE

    # ── Footer ────────────────────────────────────────────────────────────────
    y_footer = H_SLIDE - H_FOOTER
    add_rect(slide, 0, y_footer, W_SLIDE, H_FOOTER, fill_color=C_LIGHT_GRAY)

    footer_text = (
        f"Recall: {len(topk)} product(s) matched"
        f"   |   Provider: {provider}"
        f"   |   Status: {status}"
    )
    ftxb = add_textbox(slide, PAD, y_footer + 0.07, W_SLIDE - PAD * 2, 0.25,
                       footer_text, font_size=7.5, color=RGBColor(0x60,0x60,0x60))
    set_tf_margin(ftxb, 0, 0)

    return slide


def main():
    with open(JSON_PATH, encoding="utf-8") as f:
        records = json.load(f)

    prs = Presentation()
    prs.slide_width  = Inches(W_SLIDE)
    prs.slide_height = Inches(H_SLIDE)

    for record in records:
        add_slide(prs, record)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))
    print(f"Saved {len(records)} slides → {OUT_PATH}")


if __name__ == "__main__":
    main()
