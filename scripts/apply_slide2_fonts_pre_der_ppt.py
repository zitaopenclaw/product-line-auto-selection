"""
Read font sizes from slide 2 (user-edited) and apply them to slides 3-20 of the
Pre-DER Agent presentation. Matches text runs by content pattern, not shape index.
"""
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt

PPTX_PATH = Path(__file__).parent.parent / "output/pre_der_agent/presentation_pre_der_exp1.pptx"

LABELS = {"SALES INPUT", "AI EXTRACTED DESCRIPTION", "TOP RECOMMENDED PRODUCTS"}
CONF   = {"High", "Medium", "Low"}
BG_IDS = {"IDG", "DCG", "SSG"}


def classify(text, x_in):
    """Return (size_pt_or_None, bold_or_None, italic_or_None) for a text run."""
    t = text.strip()
    if not t:
        return None, None, None

    # Slide title: starts with "#<digits>  " (header, left side)
    if re.match(r'^#\d+\s', t) and x_in < 11.0:
        return 28, True, False

    # BG badge
    if t in BG_IDS:
        return 20, True, None

    # Section labels
    if t == "SALES INPUT":
        return 14, True, False
    if t == "AI EXTRACTED DESCRIPTION":
        return 14, True, False
    if t == "TOP RECOMMENDED PRODUCTS":
        return 10, True, False

    # Rank badge: exactly "#1" / "#2" / "#3"
    if re.match(r'^#\d$', t):
        return 32, True, False

    # Confidence badge
    if t in CONF:
        return 16, True, None

    # Path breadcrumb
    if " > " in t:
        return 16, False, True

    # Score line
    if t.startswith("Score:"):
        return 16, False, False

    # Footer
    if t.startswith("Recall:"):
        return 14, False, False

    # Left panel body (sales input or AI description) → clear explicit size (inherit)
    if x_in < 5.0:
        return None, None, None

    # Right panel: product name (bold text not matching above patterns)
    return 20, True, False


def apply_fonts(prs):
    updated = 0
    for slide_idx in range(2, len(prs.slides)):  # slides 3-20 (0-indexed: 2-19)
        slide = prs.slides[slide_idx]
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            x_in = shape.left / 914400  # EMU → inches
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    size_pt, bold, italic = classify(run.text, x_in)
                    run.font.size  = Pt(size_pt) if size_pt is not None else None
                    if bold   is not None: run.font.bold   = bold
                    if italic is not None: run.font.italic = italic
                    updated += 1
    return updated


def main():
    prs = Presentation(str(PPTX_PATH))
    n = apply_fonts(prs)
    prs.save(str(PPTX_PATH))
    print(f"Updated {n} text runs across slides 3-{len(prs.slides)}.")
    print(f"Saved → {PPTX_PATH}")


if __name__ == "__main__":
    main()
